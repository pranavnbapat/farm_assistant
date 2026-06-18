import csv
import io

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.routers.files import DocumentExportIn, _export_content_from_body
from app.services.document_export_service import generate_document

CONTENT = """## Crop comparison

| Crop | Rotation interval | Benefit |
| --- | --- | --- |
| Wheat | 3 years | Breaks disease cycles |
| Peas | 4 years | Adds nitrogen |
"""

def test_generate_pdf():
    document = generate_document(
        "Crop rotation",
        CONTENT,
        "pdf",
        sources=[{"title": "Crop rotation guide", "url": "https://example.com/crop-rotation"}],
    )
    assert document.filename == "Crop-rotation.pdf"
    assert document.payload.startswith(b"%PDF")
    reader = PdfReader(io.BytesIO(document.payload))
    assert len(reader.pages) == 1
    text = reader.pages[0].extract_text()
    assert "EU-FarmBook" in text
    assert "Generated" in text
    assert "Sources" in text
    assert "Crop rotation guide" in text
    assert "Page 1" in text
    assert len(PdfReader(io.BytesIO(document.payload)).pages) == 1


def test_generate_docx():
    document = generate_document("Crop rotation", CONTENT, "docx")
    parsed = Document(io.BytesIO(document.payload))
    assert parsed.paragraphs[0].text == "Crop rotation"
    assert parsed.tables[0].cell(1, 0).text == "Wheat"


def test_generate_csv():
    document = generate_document("Crop rotation", CONTENT, "csv")
    rows = list(csv.reader(io.StringIO(document.payload.decode("utf-8-sig"))))
    assert rows[0] == ["Crop", "Rotation interval", "Benefit"]
    assert rows[2][0] == "Peas"


def test_generate_xlsx():
    document = generate_document("Crop rotation", CONTENT, "xlsx")
    workbook = load_workbook(io.BytesIO(document.payload))
    sheet = workbook.active
    assert sheet["A1"].value == "Crop"
    assert sheet["A3"].value == "Peas"


def test_generate_pptx():
    document = generate_document("Crop rotation", CONTENT, "pptx")
    presentation = Presentation(io.BytesIO(document.payload))
    assert len(presentation.slides) >= 2
    assert presentation.slides[0].shapes.title.text == "Crop rotation"


def test_filename_is_sanitized():
    document = generate_document("../../Unsafe title!", "Body", "csv")
    assert document.filename == "Unsafe-title.csv"


def test_unicode_filename_is_preserved():
    document = generate_document("Précision agricole", "Body", "pdf")
    assert document.filename == "Précision-agricole.pdf"



def test_conversation_export_content_uses_transcript_rows_for_xlsx():
    body = DocumentExportIn(
        title="Conversation",
        format="xlsx",
        scope="conversation",
        messages=[
            {"role": "user", "content": "Who are you?"},
            {"role": "assistant", "content": "I am Farm Assistant."},
        ],
    )
    content = _export_content_from_body(body)
    document = generate_document(body.title, content, body.format)
    workbook = load_workbook(io.BytesIO(document.payload))
    sheet = workbook.active

    assert sheet["A1"].value == "Turn"
    assert sheet["B1"].value == "Speaker"
    assert sheet["C1"].value == "Message"
    assert sheet["B2"].value == "User"
    assert sheet["C3"].value == "I am Farm Assistant."


def test_conversation_export_content_uses_readable_transcript_for_pdf():
    body = DocumentExportIn(
        title="Conversation",
        format="pdf",
        scope="conversation",
        messages=[
            {"role": "user", "content": "Can anyone upload to EUF?"},
            {"role": "assistant", "content": "I cannot confirm that."},
        ],
    )

    content = _export_content_from_body(body)

    assert "## Conversation transcript" in content
    assert "### 1. User" in content
    assert "### 2. Farm Assistant" in content
    assert "I cannot confirm that." in content
