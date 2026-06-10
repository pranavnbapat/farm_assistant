import csv
import io
import re
from datetime import date
from dataclasses import dataclass
from html import escape
from pathlib import Path
from typing import Literal

ExportFormat = Literal["pdf", "docx", "csv", "xlsx", "pptx"]

MIME_TYPES: dict[ExportFormat, str] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "csv": "text/csv; charset=utf-8",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


@dataclass(frozen=True)
class GeneratedDocument:
    payload: bytes
    filename: str
    media_type: str


def _safe_filename(title: str, export_format: ExportFormat) -> str:
    stem = re.sub(r"[^\w.-]+", "-", title.strip(), flags=re.UNICODE).strip("-._")
    return f"{(stem or 'farm-assistant-response')[:80]}.{export_format}"


def _plain_text(markdown: str) -> str:
    text = re.sub(r"```(?:\w+)?\n?(.*?)```", r"\1", markdown, flags=re.DOTALL)
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*[-*+]\s+", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+\.\s+", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"[*_~`]+", "", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _split_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _markdown_table(markdown: str) -> list[list[str]]:
    lines = markdown.splitlines()
    for index in range(len(lines) - 1):
        header = lines[index].strip()
        separator = _split_table_row(lines[index + 1])
        if "|" not in header or not separator or not all(re.fullmatch(r":?-{3,}:?", cell) for cell in separator):
            continue
        rows = [_split_table_row(header)]
        for line in lines[index + 2:]:
            if "|" not in line:
                break
            rows.append(_split_table_row(line))
        return rows
    return []


def _tabular_rows(title: str, content: str) -> list[list[str]]:
    table = _markdown_table(content)
    if table:
        return table
    rows = [["Section", "Content"]]
    for index, paragraph in enumerate(_plain_text(content).split("\n\n")):
        if paragraph.strip():
            rows.append([title if index == 0 else "", paragraph.strip()])
    return rows


def _generate_csv(title: str, content: str) -> bytes:
    output = io.StringIO(newline="")
    csv.writer(output).writerows(_tabular_rows(title, content))
    return ("\ufeff" + output.getvalue()).encode("utf-8")


def _generate_xlsx(title: str, content: str) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Farm Assistant"
    for row in _tabular_rows(title, content):
        sheet.append(row)
    for cell in sheet[1]:
        cell.font = Font(bold=True)
    sheet.freeze_panes = "A2"
    for column in sheet.columns:
        width = min(60, max(12, max(len(str(cell.value or "")) for cell in column) + 2))
        sheet.column_dimensions[column[0].column_letter].width = width
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _generate_docx(title: str, content: str) -> bytes:
    from docx import Document

    document = Document()
    document.add_heading(title, level=0)
    table_rows = _markdown_table(content)
    if table_rows:
        table = document.add_table(rows=1, cols=max(len(row) for row in table_rows))
        table.style = "Table Grid"
        for row_index, row in enumerate(table_rows):
            cells = table.rows[0].cells if row_index == 0 else table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = value
    else:
        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            heading = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading:
                document.add_heading(heading.group(2), level=min(len(heading.group(1)), 6))
            elif re.match(r"^[-*+]\s+", line):
                document.add_paragraph(re.sub(r"^[-*+]\s+", "", line), style="List Bullet")
            else:
                document.add_paragraph(_plain_text(line))
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _generate_pdf(title: str, content: str, sources: list[dict[str, str]] | None = None) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    font_name = "Helvetica"
    font_path = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
    if font_path.exists():
        pdfmetrics.registerFont(TTFont("DejaVuSans", str(font_path)))
        font_name = "DejaVuSans"
    output = io.BytesIO()
    document = SimpleDocTemplate(
        output, pagesize=A4, rightMargin=20 * mm, leftMargin=20 * mm,
        topMargin=18 * mm, bottomMargin=20 * mm, title=title,
        author="EU-FarmBook Farm Assistant", subject="Agricultural knowledge report",
    )
    styles = getSampleStyleSheet()
    brand_style = ParagraphStyle(
        "ExportBrand", parent=styles["BodyText"], fontName=font_name,
        fontSize=9, textColor=colors.HexColor("#5B4DB7"), alignment=TA_CENTER,
    )
    title_style = ParagraphStyle(
        "ExportTitle", parent=styles["Title"], fontName=font_name,
        fontSize=22, leading=27, textColor=colors.HexColor("#162447"), alignment=TA_CENTER,
    )
    meta_style = ParagraphStyle(
        "ExportMeta", parent=styles["BodyText"], fontName=font_name,
        fontSize=8.5, textColor=colors.HexColor("#718096"), alignment=TA_CENTER,
    )
    body_style = ParagraphStyle("ExportBody", parent=styles["BodyText"], fontName=font_name, leading=16, spaceAfter=8)
    story = [
        Paragraph("EU-FarmBook&nbsp;&nbsp;|&nbsp;&nbsp;Farm Assistant", brand_style),
        Spacer(1, 4 * mm),
        Paragraph(escape(title), title_style),
        Spacer(1, 2 * mm),
        Paragraph(f"Generated {date.today().strftime('%d %B %Y')}", meta_style),
        Spacer(1, 8 * mm),
    ]
    table_rows = _markdown_table(content)
    if table_rows:
        column_widths = [174 * mm / max(1, len(table_rows[0]))] * len(table_rows[0])
        table = Table(
            [[Paragraph(escape(cell), body_style) for cell in row] for row in table_rows],
            colWidths=column_widths,
            repeatRows=1,
        )
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8EEF8")),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B8C2D8")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(table)
    else:
        for paragraph in _plain_text(content).split("\n\n"):
            story.append(Paragraph(escape(paragraph).replace("\n", "<br/>"), body_style))
    source_items = [source for source in (sources or []) if source.get("title") or source.get("project") or source.get("url") or source.get("display_url")]
    if source_items:
        source_heading = ParagraphStyle(
            "SourceHeading", parent=styles["Heading2"], fontName=font_name,
            textColor=colors.HexColor("#162447"), spaceBefore=12, spaceAfter=8,
        )
        source_style = ParagraphStyle(
            "ExportSource", parent=body_style, fontSize=9, leading=13,
            leftIndent=12, firstLineIndent=-12,
        )
        story.extend([Spacer(1, 5 * mm), Paragraph("Sources", source_heading)])
        for index, source in enumerate(source_items, start=1):
            label = source.get("title") or source.get("project") or source.get("url") or f"Source {index}"
            url = source.get("url") or source.get("display_url") or ""
            rendered_label = (
                f'<link href="{escape(url, quote=True)}" color="#3156A3">{escape(label)}</link>'
                if url else escape(label)
            )
            story.append(Paragraph(f"{index}. {rendered_label}", source_style))

    def add_footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#D9DFEA"))
        canvas.line(20 * mm, 14 * mm, 190 * mm, 14 * mm)
        canvas.setFillColor(colors.HexColor("#718096"))
        canvas.setFont("Helvetica", 8)
        canvas.drawString(20 * mm, 9 * mm, "EU-FarmBook Farm Assistant")
        canvas.drawRightString(190 * mm, 9 * mm, f"Page {doc.page}")
        canvas.restoreState()

    document.build(story, onFirstPage=add_footer, onLaterPages=add_footer)
    return output.getvalue()


def _generate_pptx(title: str, content: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Generated by Farm Assistant"
    table_rows = _markdown_table(content)
    if table_rows:
        columns = max(len(row) for row in table_rows)
        for offset in range(1, len(table_rows), 8):
            chunk = [table_rows[0], *table_rows[offset:offset + 8]]
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = title if offset == 1 else f"{title} ({offset // 8 + 1})"
            shape = slide.shapes.add_table(len(chunk), columns, Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
            table = shape.table
            for row_index, row in enumerate(chunk):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = value
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.bold = row_index == 0
    else:
        blocks = [block.strip() for block in _plain_text(content).split("\n\n") if block.strip()]
        for offset in range(0, len(blocks), 5):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = title if offset == 0 else f"{title} ({offset // 5 + 1})"
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, block in enumerate(blocks[offset:offset + 5]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = block[:900]
                paragraph.font.size = Pt(18)
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def generate_document(
    title: str,
    content: str,
    export_format: ExportFormat,
    sources: list[dict[str, str]] | None = None,
) -> GeneratedDocument:
    generators = {
        "docx": _generate_docx,
        "csv": _generate_csv,
        "xlsx": _generate_xlsx,
        "pptx": _generate_pptx,
    }
    payload = _generate_pdf(title, content, sources) if export_format == "pdf" else generators[export_format](title, content)
    return GeneratedDocument(payload, _safe_filename(title, export_format), MIME_TYPES[export_format])
